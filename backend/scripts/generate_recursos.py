"""
generate_recursos.py — Genera los recursos de presentacion del concurso:
  Recursos/portada.png · Recursos/Presentacion.pptx · Recursos/presentacion.pdf

Paleta cine (aubergine + ambar). Contenido compartido entre PPTX y PDF.
Uso: python -m backend.scripts.generate_recursos
"""
from __future__ import annotations
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parents[2] / "Recursos"
OUT.mkdir(parents=True, exist_ok=True)

PLUM = "#250f33"; PLUM2 = "#4a2168"; AMBER = "#e8a33d"; TEAL = "#1fb6a6"; INK = "#f3ecfa"; MUT = "#b9a7cf"
DEMO = "https://chanchan772.github.io/DatosAbiertosMinCultura/"
REPO = "https://github.com/chanchan772/DatosAbiertosMinCultura"

SLIDES = [
    {"t": "El problema: acceso inequitativo al cine",
     "b": ["31,7 % de las admisiones ocurren en un solo municipio (Bogotá)",
           "Top 5 municipios: más de la mitad del consumo nacional",
           "3 exhibidores concentran ~72 % del mercado",
           "De 1.122 municipios, solo ~100 tienen salas activas",
           "Cine colombiano: apenas 1,45 % de las admisiones"]},
    {"t": "Fuentes de datos",
     "b": ["SIREC (DACMI): taquilla 2026, serie diaria 2007–2026, salas activas, estrenos",
           "DIVIPOLA (datos.gov.co, obligatoria): 1.122 municipios + geolocalización",
           "DANE: proyecciones municipales por edad (bandas 15–30 / 15–45 / 15–60)",
           "Reconciliación territorial 100 % · datos anonimizados antes de procesar"]},
    {"t": "Metodología: demanda insatisfecha",
     "b": ["potencial = población objetivo × tasa de referencia (robusta, sin polos)",
           "Sin oferta → insatisfecha = potencial (brecha estructural)",
           "Subatendido → max(0, potencial − realizada)",
           "Saturado / polo → 0 (atrae público de vecinos)",
           "Distancia haversine al polo: aislamiento real vs conurbación aparente"]},
    {"t": "Resultados clave",
     "b": ["~1.025 municipios sin cine → 25,7 M de demanda potencial sin atender",
           "833 municipios de aislamiento real vs 191 de conurbación aparente",
           "Cobertura residente ≈ 54 % · demanda insatisfecha ≈ 32,8 M",
           "El problema es de DISTRIBUCIÓN territorial, no de volumen agregado"]},
    {"t": "Simulador de exhibidor",
     "b": ["captura = N_salas × rendimiento por sala del municipio",
           "Se descompone en DEMANDA NUEVA vs REDISTRIBUCIÓN de mercado",
           "Ej.: 8 salas en Bogotá → 100 % redistribución (mercado saturado)",
           "Evita presentar como crecimiento lo que le quita público a competidores"]},
    {"t": "Proyección 2027 (honesta)",
     "b": ["El sector se estabilizó en ~50 M/año (2023–2025), no crece",
           "Proyección como RANGO: 44,6 M – 55,2 M (central ~51 M)",
           "Backtest walk-forward: ~13 % de error a 1 año · R² bajo declarado",
           "2020–2021 (COVID) tratados como intervención"]},
    {"t": "Transparencia e impacto",
     "b": ["Cada cifra trae panel “¿Cómo se calcula?”: fuente, fórmula, pasos",
           "Auditado por un panel de 3 revisores; correcciones aplicadas",
           "Focaliza el fomento de la Ley 814 con evidencia cuantitativa",
           "Código abierto, reproducible y desplegado como demo pública"]},
    {"t": "Accesos",
     "b": [f"Demo en vivo:  {DEMO}",
           f"Repositorio:   {REPO}",
           "Stack: Angular 21 (frontend) + Python/FastAPI (backend) + DeepSeek",
           "Documentación técnica en docs/ · datasets en Datos/"]},
]

TITLE = "CinePredict"
SUB = "Modelo predictivo de espectadores de cine en Colombia"
RETO = "Reto 8 · Cultura y Turismo — Datos al Ecosistema 2026"
TEAM = "Grupo de Producción y Gestión de la Información · DACMI · Ministerio de las Culturas"


# ------------------------------------------------------------------ PORTADA PNG
def portada_png():
    fig = plt.figure(figsize=(12.8, 7.2), dpi=100)
    fig.patch.set_facecolor(PLUM)
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
    ax.set_facecolor(PLUM)
    # franja de "film strip"
    for i in range(14):
        ax.add_patch(plt.Rectangle((0.06 + i * 0.066, 0.12), 0.03, 0.045, color=AMBER, alpha=.85))
    ax.text(0.06, 0.72, "◐", fontsize=70, color=AMBER, transform=ax.transAxes, ha="left", va="center")
    ax.text(0.14, 0.72, TITLE, fontsize=64, color="white", weight="bold", transform=ax.transAxes, va="center")
    ax.text(0.06, 0.56, SUB, fontsize=26, color=INK, transform=ax.transAxes)
    ax.text(0.06, 0.47, RETO, fontsize=17, color=AMBER, transform=ax.transAxes)
    ax.text(0.06, 0.30, TEAM, fontsize=13, color=MUT, transform=ax.transAxes)
    ax.text(0.06, 0.22, DEMO, fontsize=13, color=TEAL, transform=ax.transAxes)
    fig.savefig(OUT / "portada.png", facecolor=PLUM)
    plt.close(fig)
    print("  [ok] portada.png")


# ------------------------------------------------------------------ PDF
def _pdf_slide(pdf, title, bullets, cover=False):
    fig = plt.figure(figsize=(13.333, 7.5), dpi=100); fig.patch.set_facecolor(PLUM if cover else "#faf9fe")
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off")
    if cover:
        ax.set_facecolor(PLUM)
        ax.text(0.07, 0.62, TITLE, fontsize=60, color="white", weight="bold", transform=ax.transAxes)
        ax.text(0.07, 0.50, SUB, fontsize=24, color=INK, transform=ax.transAxes)
        ax.text(0.07, 0.42, RETO, fontsize=16, color=AMBER, transform=ax.transAxes)
        ax.text(0.07, 0.30, TEAM, fontsize=12, color=MUT, transform=ax.transAxes)
    else:
        ax.add_patch(plt.Rectangle((0, 0.88), 1, 0.12, color=PLUM, transform=ax.transAxes))
        ax.text(0.06, 0.94, title, fontsize=26, color="white", weight="bold", transform=ax.transAxes, va="center")
        y = 0.76
        for b in bullets:
            ax.text(0.07, y, "●", fontsize=13, color=AMBER, transform=ax.transAxes, va="center")
            ax.text(0.10, y, b, fontsize=17, color="#1a1726", transform=ax.transAxes, va="center")
            y -= 0.11
    pdf.savefig(fig, facecolor=fig.get_facecolor()); plt.close(fig)


def pdf_deck():
    with PdfPages(OUT / "presentacion.pdf") as pdf:
        _pdf_slide(pdf, None, None, cover=True)
        for s in SLIDES:
            _pdf_slide(pdf, s["t"], s["b"])
    print("  [ok] presentacion.pdf")


# ------------------------------------------------------------------ PPTX
def _rgb(h): return RGBColor.from_string(h.lstrip("#").upper())

def pptx_deck():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # portada
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = _rgb(PLUM)
    def box(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = _rgb(color); return tb
    box(s, 0.8, 2.2, 11.5, 1.2, "CinePredict", 54, "FFFFFF", True)
    box(s, 0.8, 3.4, 11.5, 0.8, SUB, 24, INK)
    box(s, 0.8, 4.2, 11.5, 0.6, RETO, 16, AMBER)
    box(s, 0.8, 5.4, 11.5, 0.6, TEAM, 12, MUT)
    box(s, 0.8, 5.9, 11.5, 0.5, DEMO, 13, TEAL)
    for sl in SLIDES:
        s = prs.slides.add_slide(blank)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = _rgb("FAF9FE")
        band = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        band.fill.solid(); band.fill.fore_color.rgb = _rgb(PLUM); band.line.fill.background()
        tf = band.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "  " + sl["t"]; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = _rgb("FFFFFF")
        tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5)); body = tb.text_frame; body.word_wrap = True
        for i, b in enumerate(sl["b"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.space_after = Pt(14)
            rb = p.add_run(); rb.text = "●  "; rb.font.color.rgb = _rgb(AMBER); rb.font.size = Pt(18)
            rt = p.add_run(); rt.text = b; rt.font.size = Pt(18); rt.font.color.rgb = _rgb("1A1726")
    prs.save(OUT / "Presentacion.pptx")
    print("  [ok] Presentacion.pptx")


if __name__ == "__main__":
    print("Generando recursos en", OUT)
    portada_png(); pdf_deck(); pptx_deck(); print("Listo.")
